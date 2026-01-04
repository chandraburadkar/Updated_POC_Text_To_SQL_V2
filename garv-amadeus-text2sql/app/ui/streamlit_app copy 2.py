# app/ui/streamlit_app.py
from __future__ import annotations

import time
from io import BytesIO
from typing import Any, Dict, List, Optional

import pandas as pd
import requests
import streamlit as st

# Optional exports (PDF/PPT). UI will gracefully disable if missing.
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
except Exception:  # pragma: no cover
    canvas = None  # type: ignore

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


# =============================
# CONFIG
# =============================
API_BASE = "http://127.0.0.1:8000"
HEALTH_ENDPOINT = "/api/health"
TEXT2SQL_ENDPOINT = "/api/text2sql"

APP_TITLE = "GARV – Airport Operations Intelligence"
APP_SUBTITLE = "Enterprise conversational analytics (Text-to-SQL) for Airport Ops"


# =============================
# PAGE SETUP
# =============================
st.set_page_config(
    page_title="GARV – Airport Ops Intelligence",
    page_icon="✈️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# =============================
# CSS (enterprise, sober)
# =============================
st.markdown(
    """
<style>
.block-container { padding-top: 0.75rem !important; max-width: 1220px; }
header[data-testid="stHeader"] { display: none; }
section[data-testid="stSidebar"] { border-right: 1px solid #e5e7eb; }
.small-muted { color: #6b7280; font-size: 12px; }
.hint { color: #6b7280; font-size: 13px; }
.badge { display: inline-block; padding: 2px 10px; border-radius: 999px; font-size: 12px; border: 1px solid #e5e7eb; }
.badge-high { background: #ecfdf5; border-color: #a7f3d0; color: #065f46; }
.badge-med { background: #eff6ff; border-color: #bfdbfe; color: #1e3a8a; }
.badge-low { background: #fff7ed; border-color: #fed7aa; color: #9a3412; }
.card { border: 1px solid #e5e7eb; border-radius: 14px; padding: 16px; background: #fff; }
.section-title { font-size: 14px; color: #111827; font-weight: 700; margin-bottom: 6px; }
.hr { height: 1px; background: #e5e7eb; margin: 14px 0; }
.right-header { display:flex; align-items:center; justify-content:space-between; gap:12px; }
.right-header h2 { margin: 0; }
.mono { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono", "Courier New", monospace; }
</style>
""",
    unsafe_allow_html=True,
)


# =============================
# HEALTH CHECK
# =============================
def is_api_alive() -> bool:
    try:
        r = requests.get(API_BASE + HEALTH_ENDPOINT, timeout=1.5)
        return r.status_code == 200
    except Exception:
        return False


# =============================
# SESSION STATE
# =============================
if "chats" not in st.session_state:
    st.session_state.chats = []
if "active_chat" not in st.session_state:
    st.session_state.active_chat = None
if "saved_insights" not in st.session_state:
    st.session_state.saved_insights = []
if "show_diagnostics" not in st.session_state:
    st.session_state.show_diagnostics = False


def _now_ms() -> int:
    return int(time.time() * 1000)


def new_chat() -> None:
    cid = f"chat-{_now_ms()}"
    st.session_state.chats.insert(
        0,
        {
            "id": cid,
            "title": "New analysis",
            "messages": [
                {
                    "role": "assistant",
                    "type": "text",
                    "content": "Ask an Airport Ops question (e.g., “Top 5 airports by avg security wait time in last 30 days”).",
                }
            ],
            "chat_history": [],
            "memory_entities": {},
            # per-chat viz edits
            "viz_overrides": {},  # message_idx -> override dict
        },
    )
    st.session_state.active_chat = cid


def get_chat() -> Dict[str, Any]:
    if not st.session_state.active_chat:
        new_chat()
    for c in st.session_state.chats:
        if c["id"] == st.session_state.active_chat:
            c.setdefault("chat_history", [])
            c.setdefault("memory_entities", {})
            c.setdefault("messages", [])
            c.setdefault("viz_overrides", {})
            return c
    new_chat()
    return get_chat()


def _set_active(cid: str) -> None:
    st.session_state.active_chat = cid


# =============================
# EXPORT HELPERS
# =============================
def df_to_excel_bytes(df: pd.DataFrame) -> bytes:
    buf = BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Results")
    return buf.getvalue()


def _wrap_text(text: str, width: int) -> List[str]:
    words = (text or "").split()
    out: List[str] = []
    line: List[str] = []
    for w in words:
        if len(" ".join(line + [w])) <= width:
            line.append(w)
        else:
            out.append(" ".join(line))
            line = [w]
    if line:
        out.append(" ".join(line))
    return out


def build_pdf_bytes(payload: Dict[str, Any]) -> Optional[bytes]:
    """
    Lightweight PDF export (summary + SQL + small table preview).
    Requires reportlab installed.
    """
    if canvas is None:
        return None

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter

    y = height - 50
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, y, "GARV – Airport Operations Intelligence")
    y -= 22
    c.setFont("Helvetica", 10)
    c.drawString(50, y, "Exported insight")
    y -= 18

    answer_card = payload.get("answer_card") or {}
    answer = (answer_card.get("answer") or "").strip()
    sql = (payload.get("final_sql") or "").strip()

    # Summary
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Executive Insight")
    y -= 16
    c.setFont("Helvetica", 10)
    for line in (answer[:900] or "").split("\n"):
        for chunk in _wrap_text(line, 95):
            c.drawString(50, y, chunk)
            y -= 12
            if y < 80:
                c.showPage()
                y = height - 50

    y -= 8

    # SQL
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Query Logic (SQL)")
    y -= 16
    c.setFont("Helvetica", 9)
    for line in _wrap_text(sql, 110):
        c.drawString(50, y, line)
        y -= 11
        if y < 80:
            c.showPage()
            y = height - 50

    # Table preview (first 15 rows)
    rows = payload.get("rows") or []
    if isinstance(rows, list) and rows:
        y -= 6
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, y, "Result Preview")
        y -= 16
        c.setFont("Helvetica", 9)
        df = pd.DataFrame(rows).head(15)
        preview = df.to_string(index=False)
        for line in preview.split("\n"):
            c.drawString(50, y, line[:120])
            y -= 11
            if y < 80:
                c.showPage()
                y = height - 50

    c.showPage()
    c.save()
    return buf.getvalue()


def build_pptx_bytes(payload: Dict[str, Any]) -> Optional[bytes]:
    """
    Lightweight PPT export (title + insight + SQL + small table preview).
    Requires python-pptx installed.
    """
    if Presentation is None:
        return None

    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    tf = title.text_frame
    tf.text = "GARV – Airport Operations Intelligence"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True

    answer_card = payload.get("answer_card") or {}
    insight = (answer_card.get("answer") or "").strip() or "Executive insight not available."

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.2), Inches(2.2))
    tf2 = box.text_frame
    tf2.text = "Executive Insight"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    p = tf2.add_paragraph()
    p.text = insight[:1400]
    p.font.size = Pt(12)

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    t2 = slide2.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    tf = t2.text_frame
    tf.text = "Query Logic & Results"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True

    sql = (payload.get("final_sql") or "").strip()
    sql_box = slide2.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12.2), Inches(2.2))
    tf3 = sql_box.text_frame
    tf3.text = "SQL"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    p = tf3.add_paragraph()
    p.text = sql[:1600]
    p.font.size = Pt(10)

    rows = payload.get("rows") or []
    if isinstance(rows, list) and rows:
        df = pd.DataFrame(rows).head(12)
        preview = df.to_string(index=False)
        prev_box = slide2.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(12.2), Inches(3.2))
        tf4 = prev_box.text_frame
        tf4.text = "Result Preview"
        tf4.paragraphs[0].font.size = Pt(14)
        tf4.paragraphs[0].font.bold = True
        p2 = tf4.add_paragraph()
        p2.text = preview[:1800]
        p2.font.size = Pt(9)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =============================
# TRACE GRAPH (post-run)
# =============================
def trace_to_dot(trace: List[Dict[str, Any]]) -> str:
    """
    Visualize executed steps AFTER completion.
    (True realtime streaming requires server streaming/WebSockets; UI-only can’t stream step-by-step.)
    """
    nodes: List[str] = []
    edges: List[str] = []
    for i, t in enumerate(trace or []):
        step = (t.get("step") or f"step_{i}").replace('"', "'")
        ok = bool(t.get("ok", False))
        label = f"{i+1}. {step}\\n{'OK' if ok else 'FAIL'}"
        color = "palegreen" if ok else "mistyrose"
        nodes.append(f'n{i} [label="{label}", style="filled", fillcolor="{color}"];')
        if i > 0:
            edges.append(f"n{i-1} -> n{i};")
    dot = (
        "digraph G {\n"
        "rankdir=LR;\n"
        "node [shape=box, fontname=Helvetica];\n"
        + "\n".join(nodes)
        + "\n"
        + "\n".join(edges)
        + "\n}"
    )
    return dot


# =============================
# CHART RENDER (auto + editable)
# =============================
def _infer_columns(df: pd.DataFrame) -> Dict[str, List[str]]:
    cols = list(df.columns)
    numeric = [c for c in cols if pd.api.types.is_numeric_dtype(df[c])]
    datetime_cols = [c for c in cols if pd.api.types.is_datetime64_any_dtype(df[c])]
    categorical = [c for c in cols if c not in numeric and c not in datetime_cols]
    return {"cols": cols, "numeric": numeric, "datetime": datetime_cols, "categorical": categorical}


def render_chart(
    df: pd.DataFrame,
    chart_spec: Optional[Dict[str, Any]],
    override: Optional[Dict[str, Any]] = None,
) -> None:
    """
    Editable chart rendering.
    """
    if df is None or df.empty:
        st.info("No data to visualize.")
        return

    meta = _infer_columns(df)
    cols = meta["cols"]

    spec = dict(chart_spec or {})
    if override:
        spec.update({k: v for k, v in override.items() if v is not None})

    ctype = (spec.get("type") or spec.get("plot") or "bar").strip().lower()
    x = spec.get("x")
    y = spec.get("y")

    if not x:
        x = meta["datetime"][0] if meta["datetime"] else (meta["categorical"][0] if meta["categorical"] else cols[0])
    if not y:
        y = meta["numeric"][0] if meta["numeric"] else None

    if ctype in ("line", "timeseries"):
        d = df.copy()
        if x in d.columns:
            try:
                d[x] = pd.to_datetime(d[x], errors="ignore")
                d = d.sort_values(by=x)
                d = d.set_index(x)
            except Exception:
                pass
        if y and y in d.columns:
            st.line_chart(d[y])
        else:
            st.line_chart(d.select_dtypes(include="number"))
    elif ctype in ("bar", "column"):
        if x in df.columns and y and y in df.columns:
            d = df[[x, y]].copy()
            try:
                d2 = d.groupby(x, as_index=True)[y].mean().sort_values(ascending=False).head(20)
                st.bar_chart(d2)
            except Exception:
                st.bar_chart(d.set_index(x)[y])
        else:
            st.bar_chart(df.select_dtypes(include="number"))
    elif ctype == "scatter":
        if x in df.columns and y and y in df.columns:
            try:
                st.scatter_chart(df, x=x, y=y)
            except Exception:
                st.dataframe(df[[x, y]])
        else:
            st.dataframe(df)
    else:
        st.dataframe(df, use_container_width=True, height=320)


# =============================
# SIDEBAR (enterprise)
# =============================
with st.sidebar:
    st.markdown(f"### {APP_TITLE}")
    st.caption(APP_SUBTITLE)

    if st.button("➕ New analysis", use_container_width=True, key="btn_new_chat"):
        new_chat()
        st.rerun()

    st.divider()

    st.markdown("**Recent Analyses**")
    if not st.session_state.chats:
        st.caption("No analyses yet.")
    else:
        for c in st.session_state.chats[:12]:
            title = c.get("title") or "Untitled"
            if st.button(title, use_container_width=True, key=f"chat_select_{c['id']}"):
                _set_active(c["id"])
                st.rerun()

    st.divider()

    st.markdown("**Saved Insights**")
    if not st.session_state.saved_insights:
        st.caption("Save an insight from the right panel to view it here.")
    else:
        for s in st.session_state.saved_insights[:12]:
            if st.button(s["title"], use_container_width=True, key=f"saved_{s['id']}"):
                chat = get_chat()
                chat["messages"].append(
                    {"role": "assistant", "type": "answer", "payload": s["payload"], "saved": True}
                )
                st.rerun()

    st.divider()
    connected = is_api_alive()
    st.caption("Status: 🟢 Connected" if connected else "Status: 🔴 API Offline")


# =============================
# MAIN (stacked, ChatGPT-like history)
# =============================
chat = get_chat()

st.markdown(
    f"""
<div class="right-header">
  <h2>{chat.get("title","New analysis")}</h2>
</div>
""",
    unsafe_allow_html=True,
)
st.caption("GARV returns an executive insight, SQL, results, editable visualization, and diagnostics.")


# =============================
# RENDER MESSAGES
# =============================
for idx, m in enumerate(chat["messages"]):
    role = m.get("role", "assistant")
    mtype = m.get("type", "text")
    payload = m.get("payload", {})

    with st.chat_message(role):
        if mtype == "text":
            st.markdown(m.get("content", ""))
            continue

        if mtype != "answer":
            st.markdown(m.get("content", ""))
            continue

        out = payload if isinstance(payload, dict) else {}
        answer_card = out.get("answer_card") or {}
        final_sql = out.get("final_sql") or ""
        rows = out.get("rows") or []
        df = pd.DataFrame(rows) if isinstance(rows, list) and rows else pd.DataFrame()
        chart_spec = out.get("chart_spec") or out.get("chart") or None
        debug = out.get("debug") or {}
        trace = (debug.get("trace") or []) if isinstance(debug, dict) else []

        # 1) Question
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Question</div>', unsafe_allow_html=True)
        st.markdown(f"“{out.get('question') or out.get('user_question') or ''}”")
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # 2) Executive Insight
        conf = (answer_card.get("confidence") or "MEDIUM").upper()
        badge_class = "badge-med"
        if conf == "HIGH":
            badge_class = "badge-high"
        elif conf == "LOW":
            badge_class = "badge-low"

        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown(
            f"""
<div style="display:flex;align-items:center;justify-content:space-between;gap:10px;">
  <div class="section-title">Executive Insight</div>
  <span class="badge {badge_class}">Confidence: {conf}</span>
</div>
""",
            unsafe_allow_html=True,
        )

        insight = (answer_card.get("answer") or "").strip()
        if not insight:
            exp = out.get("explanation") or {}
            if isinstance(exp, dict):
                insight = (exp.get("summary") or "").strip()
        st.markdown(insight or "Insight not available.")

        scope = answer_card.get("scope") or {}
        entities = scope.get("entities") or out.get("entities") or {}
        tables = scope.get("tables") or out.get("retrieved_tables") or []

        with st.expander("Details (scope & assumptions)", expanded=False):
            st.markdown("**Scope**")
            st.write(
                {
                    "intent": out.get("intent"),
                    "entities": entities,
                    "tables_used": tables,
                }
            )
            assumptions = out.get("assumptions") or []
            if assumptions:
                st.markdown("**Assumptions**")
                for a in assumptions:
                    st.markdown(f"- {a}")

        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # 3) Query Logic (SQL)
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Query Logic (SQL)</div>', unsafe_allow_html=True)
        with st.expander("View SQL", expanded=False):
            st.code(final_sql, language="sql")
            st.download_button(
                "Download SQL",
                data=final_sql.encode("utf-8"),
                file_name=f"{chat['id']}_sql_{idx}.sql",
                mime="text/plain",
                key=f"dl_sql_{chat['id']}_{idx}",
            )
        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # 4) Result Data
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Result Data</div>', unsafe_allow_html=True)

        if df.empty:
            st.info("No rows returned.")
        else:
            st.caption(f"Rows returned: {len(df)}")
            st.dataframe(df, use_container_width=True, height=320)

            c1, c2, c3 = st.columns([1, 1, 2])
            with c1:
                st.download_button(
                    "⬇ Download CSV",
                    data=df.to_csv(index=False).encode("utf-8"),
                    file_name=f"{chat['id']}_results_{idx}.csv",
                    mime="text/csv",
                    key=f"dl_csv_{chat['id']}_{idx}",
                )
            with c2:
                st.download_button(
                    "⬇ Download Excel",
                    data=df_to_excel_bytes(df),
                    file_name=f"{chat['id']}_results_{idx}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"dl_xlsx_{chat['id']}_{idx}",
                )
            with c3:
                st.caption("Tip: Use Export to PDF/PPT for client-ready sharing.")

        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # 5) Visual Analysis (editable)
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Visual Analysis</div>', unsafe_allow_html=True)
        st.caption("Auto-generated visualization with an edit option (client-friendly).")

        overrides_for_chat = chat.get("viz_overrides", {})
        current_override = overrides_for_chat.get(str(idx), {}) if isinstance(overrides_for_chat, dict) else {}

        if df.empty:
            st.info("No data to visualize.")
        else:
            render_chart(df, chart_spec, override=current_override)

            with st.expander("Edit visualization", expanded=False):
                meta = _infer_columns(df)
                cols = meta["cols"]
                numeric = meta["numeric"]
                datetime_cols = meta["datetime"]
                categorical = meta["categorical"]

                default_type = (current_override.get("type") or (chart_spec or {}).get("type") or "bar")
                chart_type = st.selectbox(
                    "Chart type",
                    ["bar", "line", "scatter", "table"],
                    index=["bar", "line", "scatter", "table"].index(default_type)
                    if default_type in ["bar", "line", "scatter", "table"]
                    else 0,
                    key=f"viz_type_{chat['id']}_{idx}",
                )

                default_x = (
                    current_override.get("x")
                    or (chart_spec or {}).get("x")
                    or (datetime_cols[0] if datetime_cols else (categorical[0] if categorical else cols[0]))
                )
                default_y = (
                    current_override.get("y")
                    or (chart_spec or {}).get("y")
                    or (numeric[0] if numeric else None)
                )

                x = st.selectbox(
                    "X axis",
                    cols,
                    index=cols.index(default_x) if default_x in cols else 0,
                    key=f"viz_x_{chat['id']}_{idx}",
                )

                y: Optional[str] = None
                if numeric:
                    y = st.selectbox(
                        "Y axis",
                        numeric,
                        index=numeric.index(default_y) if default_y in numeric else 0,
                        key=f"viz_y_{chat['id']}_{idx}",
                    )
                else:
                    st.caption("No numeric column detected for Y axis.")

                if st.button("Apply changes", use_container_width=True, key=f"apply_viz_{chat['id']}_{idx}"):
                    chat.setdefault("viz_overrides", {})
                    chat["viz_overrides"][str(idx)] = {"type": chart_type, "x": x, "y": y}
                    st.rerun()

        st.markdown("</div>", unsafe_allow_html=True)

        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # 6) Actions (Save / Export)
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Actions</div>', unsafe_allow_html=True)

        a1, a2, a3 = st.columns([1, 1, 1])
        with a1:
            if st.button("💾 Save insight", use_container_width=True, key=f"save_{chat['id']}_{idx}"):
                sid = f"saved-{_now_ms()}"
                title = chat.get("title") or "Saved insight"
                st.session_state.saved_insights.insert(
                    0, {"id": sid, "title": title, "payload": out, "ts": int(time.time())}
                )
                st.success("Saved (see sidebar → Saved Insights).")
        with a2:
            pdf_bytes = build_pdf_bytes(out)
            st.download_button(
                "⬇ Export to PDF",
                data=pdf_bytes if pdf_bytes else b"",
                file_name=f"{chat['id']}_export_{idx}.pdf",
                mime="application/pdf",
                disabled=(pdf_bytes is None),
                use_container_width=True,
                key=f"pdf_{chat['id']}_{idx}",
            )
            if pdf_bytes is None:
                st.caption("Install reportlab to enable PDF export.")
        with a3:
            pptx_bytes = build_pptx_bytes(out)
            st.download_button(
                "⬇ Export to PPT",
                data=pptx_bytes if pptx_bytes else b"",
                file_name=f"{chat['id']}_export_{idx}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                disabled=(pptx_bytes is None),
                use_container_width=True,
                key=f"ppt_{chat['id']}_{idx}",
            )
            if pptx_bytes is None:
                st.caption("Install python-pptx to enable PPT export.")

        st.markdown("</div>", unsafe_allow_html=True)

        # Diagnostics graph (kept, not changed; shown as graph)
        if trace:
            with st.expander("Admin / Diagnostics (Execution Diagnostics)", expanded=False):
                st.caption("This shows which agent/step executed. Keep it collapsed during client demo.")
                st.graphviz_chart(trace_to_dot(trace))
                st.json({"trace": trace})


# =============================
# INPUT + RUN
# =============================
prompt = st.chat_input("Ask a question… e.g. Top 5 airports by avg security wait time in last 30 days")

if prompt:
    # Update right-side title (keep as you asked)
    if chat["title"] == "New analysis":
        chat["title"] = prompt[:34] + ("…" if len(prompt) > 34 else "")

    # Append user message
    chat["messages"].append({"role": "user", "type": "text", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    # Assistant response
    with st.chat_message("assistant"):
        if not is_api_alive():
            msg = "❌ API is Offline. Start FastAPI first, then retry."
            st.markdown(msg)
            chat["messages"].append({"role": "assistant", "type": "text", "content": msg})
        else:
            # Client-friendly progress (UI-only; true realtime agent streaming needs backend streaming)
            status = st.status("Running analysis…", expanded=True)
            status.update(label="Understanding request…", state="running")
            time.sleep(0.15)
            status.update(label="Grounding to airport ops schema…", state="running")
            time.sleep(0.15)
            status.update(label="Generating SQL & validating…", state="running")

            try:
                req_payload = {
                    "question": prompt,
                    "top_k_schema": 5,
                    "return_rows": 50,
                    "enable_viz": True,
                    "chat_history": chat.get("chat_history", []),
                    "memory_entities": chat.get("memory_entities", {}),
                }

                r = requests.post(API_BASE + TEXT2SQL_ENDPOINT, json=req_payload, timeout=180)
                out = r.json()

                # Update memory (if your API returns them)
                chat["chat_history"] = out.get("chat_history", chat.get("chat_history", []))
                chat["memory_entities"] = out.get("memory_entities", chat.get("memory_entities", {}))

                if out.get("ok"):
                    status.update(label="Finalizing executive insight…", state="running")
                    time.sleep(0.12)
                    status.update(label="Completed", state="complete")

                    # Ensure question is present for stacked renderer
                    out["question"] = prompt

                    # Persist the full output so renderer can show answer_card, chart_spec, debug.trace
                    chat["messages"].append({"role": "assistant", "type": "answer", "payload": out})
                    st.rerun()
                else:
                    status.update(label="Needs input", state="error")
                    stage = out.get("stage", "")
                    msg = out.get("message", "Unknown error")
                    err = f"🤔 {msg}" if stage == "clarification" else f"❌ {msg}"
                    st.markdown(err)
                    chat["messages"].append({"role": "assistant", "type": "text", "content": err})

            except Exception as e:
                status.update(label="Failed", state="error")
                err = f"❌ API error: {e}"
                st.markdown(err)
                chat["messages"].append({"role": "assistant", "type": "text", "content": err})