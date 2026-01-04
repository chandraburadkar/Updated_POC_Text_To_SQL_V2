# app/ui/streamlit_app.py
from __future__ import annotations

import time
from io import BytesIO
from typing import Any, Dict, List, Optional

import pandas as pd
import requests
import streamlit as st

# Optional exports (PDF/PPT). UI will gracefully disable if missing.
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
except Exception:  # pragma: no cover
    canvas = None  # type: ignore

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


# =============================
# CONFIG
# =============================
API_BASE = "http://127.0.0.1:8000"
HEALTH_ENDPOINT = "/api/health"
TEXT2SQL_ENDPOINT = "/api/text2sql"

APP_TITLE = "GARV – Airport Operations Intelligence"
APP_SUBTITLE = "Enterprise conversational analytics (Text-to-SQL) for Airport Ops"


# =============================
# PAGE SETUP
# =============================
st.set_page_config(
    page_title="GARV – Airport Ops Intelligence",
    page_icon="✈️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# =============================
# CSS (enterprise + clean suggestions)
# =============================
st.markdown(
    """
<style>
.block-container { padding-top: 0.65rem !important; max-width: 1220px; }
header[data-testid="stHeader"] { display: none; }
section[data-testid="stSidebar"] { border-right: 1px solid #e5e7eb; }

.small-muted { color: #6b7280; font-size: 12px; }
.hint { color: #6b7280; font-size: 13px; }

.badge { display: inline-block; padding: 2px 10px; border-radius: 999px; font-size: 12px; border: 1px solid #e5e7eb; }
.badge-high { background: #ecfdf5; border-color: #a7f3d0; color: #065f46; }
.badge-med  { background: #eff6ff; border-color: #bfdbfe; color: #1e3a8a; }
.badge-low  { background: #fff7ed; border-color: #fed7aa; color: #9a3412; }

.card { border: 1px solid #e5e7eb; border-radius: 14px; padding: 16px; background: #fff; }
.section-title { font-size: 14px; color: #111827; font-weight: 700; margin-bottom: 6px; }
.hr { height: 1px; background: #e5e7eb; margin: 14px 0; }

.right-header { display:flex; align-items:center; justify-content:space-between; gap:12px; }
.right-header h2 { margin: 0; }

.mono { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono", "Courier New", monospace; }

/* Tighten user chat bubble spacing */
div[data-testid="stChatMessage"] { padding-top: 0.15rem; padding-bottom: 0.15rem; }

/* Suggestion buttons styled as chips */
div[data-testid="stButton"] > button {
  border-radius: 999px !important;
  padding: 6px 10px !important;
  font-size: 12px !important;
  border: 1px solid #e5e7eb !important;
  background: #f9fafb !important;

  /* ✅ allow wrapping inside the button so it doesn't overlap */
  white-space: normal !important;
  word-break: break-word !important;
  text-align: left !important;
  line-height: 1.25 !important;
  height: auto !important;
}
div[data-testid="stButton"] > button:hover {
  background: #f3f4f6 !important;
  border-color: #d1d5db !important;
}
</style>
""",
    unsafe_allow_html=True,
)


# =============================
# HEALTH CHECK
# =============================
def is_api_alive() -> bool:
    try:
        r = requests.get(API_BASE + HEALTH_ENDPOINT, timeout=1.5)
        return r.status_code == 200
    except Exception:
        return False


# =============================
# SESSION STATE
# =============================
if "chats" not in st.session_state:
    st.session_state.chats = []
if "active_chat" not in st.session_state:
    st.session_state.active_chat = None
if "show_diagnostics" not in st.session_state:
    st.session_state.show_diagnostics = False


def _now_ms() -> int:
    return int(time.time() * 1000)


def new_chat() -> None:
    """Create a new empty analysis."""
    cid = f"chat-{_now_ms()}"
    st.session_state.chats.insert(
        0,
        {
            "id": cid,
            "title": "New analysis",
            "messages": [],
            "chat_history": [],
            "memory_entities": {},
            "viz_overrides": {},  # message_idx -> override dict
        },
    )
    st.session_state.active_chat = cid


def get_chat() -> Dict[str, Any]:
    if not st.session_state.active_chat:
        new_chat()
    for c in st.session_state.chats:
        if c["id"] == st.session_state.active_chat:
            c.setdefault("chat_history", [])
            c.setdefault("memory_entities", {})
            c.setdefault("messages", [])
            c.setdefault("viz_overrides", {})
            return c
    new_chat()
    return get_chat()


def _set_active(cid: str) -> None:
    st.session_state.active_chat = cid


# =============================
# EXPORT HELPERS
# =============================
def df_to_excel_bytes(df: pd.DataFrame) -> Optional[bytes]:
    """
    Excel export (optional).
    If openpyxl isn't installed, return None and disable the button.
    """
    try:
        import openpyxl  # noqa: F401
    except Exception:
        return None

    buf = BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Results")
    return buf.getvalue()


def _wrap_text(text: str, width: int) -> List[str]:
    words = (text or "").split()
    out: List[str] = []
    line: List[str] = []
    for w in words:
        if len(" ".join(line + [w])) <= width:
            line.append(w)
        else:
            out.append(" ".join(line))
            line = [w]
    if line:
        out.append(" ".join(line))
    return out


def build_pdf_bytes(payload: Dict[str, Any]) -> Optional[bytes]:
    """Lightweight PDF export (summary + SQL + small table preview). Requires reportlab."""
    if canvas is None:
        return None

    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter

    y = height - 50
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, y, "GARV – Airport Operations Intelligence")
    y -= 22
    c.setFont("Helvetica", 10)
    c.drawString(50, y, "Exported insight")
    y -= 18

    answer_card = payload.get("answer_card") or {}
    answer = (answer_card.get("answer") or "").strip()
    sql = (payload.get("final_sql") or "").strip()

    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Executive Insight")
    y -= 16
    c.setFont("Helvetica", 10)
    for line in (answer[:900] or "").split("\n"):
        for chunk in _wrap_text(line, 95):
            c.drawString(50, y, chunk)
            y -= 12
            if y < 80:
                c.showPage()
                y = height - 50

    y -= 8
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "How we calculated it")
    y -= 16
    c.setFont("Helvetica", 9)
    for line in _wrap_text(sql, 110):
        c.drawString(50, y, line)
        y -= 11
        if y < 80:
            c.showPage()
            y = height - 50

    rows = payload.get("rows") or []
    if isinstance(rows, list) and rows:
        y -= 6
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, y, "Detailed results (preview)")
        y -= 16
        c.setFont("Helvetica", 9)
        df = pd.DataFrame(rows).head(15)
        preview = df.to_string(index=False)
        for line in preview.split("\n"):
            c.drawString(50, y, line[:120])
            y -= 11
            if y < 80:
                c.showPage()
                y = height - 50

    c.showPage()
    c.save()
    return buf.getvalue()


def build_pptx_bytes(payload: Dict[str, Any]) -> Optional[bytes]:
    """Lightweight PPT export. Requires python-pptx."""
    if Presentation is None:
        return None

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    tf = title.text_frame
    tf.text = "GARV – Airport Operations Intelligence"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True

    answer_card = payload.get("answer_card") or {}
    insight = (answer_card.get("answer") or "").strip() or "Executive insight not available."

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.2), Inches(2.2))
    tf2 = box.text_frame
    tf2.text = "Executive Insight"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    p = tf2.add_paragraph()
    p.text = insight[:1400]
    p.font.size = Pt(12)

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    t2 = slide2.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    tf = t2.text_frame
    tf.text = "How we calculated it + Detailed results"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True

    sql = (payload.get("final_sql") or "").strip()
    sql_box = slide2.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12.2), Inches(2.2))
    tf3 = sql_box.text_frame
    tf3.text = "How we calculated it (SQL)"
    tf3.paragraphs[0].font.size = Pt(14)
    tf3.paragraphs[0].font.bold = True
    p = tf3.add_paragraph()
    p.text = sql[:1600]
    p.font.size = Pt(10)

    rows = payload.get("rows") or []
    if isinstance(rows, list) and rows:
        df = pd.DataFrame(rows).head(12)
        preview = df.to_string(index=False)
        prev_box = slide2.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(12.2), Inches(3.2))
        tf4 = prev_box.text_frame
        tf4.text = "Detailed results (preview)"
        tf4.paragraphs[0].font.size = Pt(14)
        tf4.paragraphs[0].font.bold = True
        p2 = tf4.add_paragraph()
        p2.text = preview[:1800]
        p2.font.size = Pt(9)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =============================
# TRACE GRAPH (post-run)
# =============================
def trace_to_dot(trace: List[Dict[str, Any]]) -> str:
    nodes: List[str] = []
    edges: List[str] = []
    for i, t in enumerate(trace or []):
        step = (t.get("step") or f"step_{i}").replace('"', "'")
        ok = bool(t.get("ok", False))
        label = f"{i+1}. {step}\\n{'OK' if ok else 'FAIL'}"
        color = "palegreen" if ok else "mistyrose"
        nodes.append(f'n{i} [label="{label}", style="filled", fillcolor="{color}"];')
        if i > 0:
            edges.append(f"n{i-1} -> n{i};")
    dot = (
        "digraph G {\n"
        "rankdir=LR;\n"
        "node [shape=box, fontname=Helvetica];\n"
        + "\n".join(nodes)
        + "\n"
        + "\n".join(edges)
        + "\n}"
    )
    return dot


# =============================
# CHART RENDER (auto + editable)
# =============================
def _infer_columns(df: pd.DataFrame) -> Dict[str, List[str]]:
    cols = list(df.columns)
    numeric = [c for c in cols if pd.api.types.is_numeric_dtype(df[c])]
    datetime_cols = [c for c in cols if pd.api.types.is_datetime64_any_dtype(df[c])]
    categorical = [c for c in cols if c not in numeric and c not in datetime_cols]
    return {"cols": cols, "numeric": numeric, "datetime": datetime_cols, "categorical": categorical}


def render_chart(
    df: pd.DataFrame,
    chart_spec: Optional[Dict[str, Any]],
    override: Optional[Dict[str, Any]] = None,
) -> None:
    if df is None or df.empty:
        st.info("No data to visualize.")
        return

    meta = _infer_columns(df)
    cols = meta["cols"]

    spec = dict(chart_spec or {})
    if override:
        spec.update({k: v for k, v in override.items() if v is not None})

    ctype = (spec.get("type") or spec.get("plot") or "bar").strip().lower()
    x = spec.get("x")
    y = spec.get("y")

    if not x:
        x = meta["datetime"][0] if meta["datetime"] else (meta["categorical"][0] if meta["categorical"] else cols[0])
    if not y:
        y = meta["numeric"][0] if meta["numeric"] else None

    if ctype in ("line", "timeseries"):
        d = df.copy()
        if x in d.columns:
            try:
                d[x] = pd.to_datetime(d[x], errors="ignore")
                d = d.sort_values(by=x)
                d = d.set_index(x)
            except Exception:
                pass
        if y and y in d.columns:
            st.line_chart(d[y])
        else:
            st.line_chart(d.select_dtypes(include="number"))
    elif ctype in ("bar", "column"):
        if x in df.columns and y and y in df.columns:
            d = df[[x, y]].copy()
            try:
                d2 = d.groupby(x, as_index=True)[y].mean().sort_values(ascending=False).head(20)
                st.bar_chart(d2)
            except Exception:
                st.bar_chart(d.set_index(x)[y])
        else:
            st.bar_chart(df.select_dtypes(include="number"))
    elif ctype == "scatter":
        if x in df.columns and y and y in df.columns:
            try:
                st.scatter_chart(df, x=x, y=y)
            except Exception:
                st.dataframe(df[[x, y]])
        else:
            st.dataframe(df)
    else:
        st.dataframe(df, use_container_width=True, height=320)


# =============================
# SUGGESTIONS (no navigation)
# =============================
def _latest_suggestions_from_chat(chat: Dict[str, Any], max_items: int = 3) -> List[str]:
    """Most recent assistant payload.suggested_questions."""
    for m in reversed(chat.get("messages", [])):
        if m.get("role") == "assistant" and m.get("type") == "answer":
            out = m.get("payload") or {}
            sugg = out.get("suggested_questions") or []
            if isinstance(sugg, list):
                cleaned: List[str] = []
                for s in sugg:
                    s = str(s).strip()
                    if s and s.lower() not in [x.lower() for x in cleaned]:
                        cleaned.append(s)
                return cleaned[:max_items]
    return []


def render_suggestions_inline(chat: Dict[str, Any]) -> None:
    """
    Renders suggestions as Streamlit buttons above the chat input.
    Clicking does NOT navigate; it sets session_state and reruns.
    """
    suggestions = _latest_suggestions_from_chat(chat, max_items=3)
    if not suggestions:
        return

    st.markdown("<div class='card' style='padding:10px 12px; margin-top:8px;'>", unsafe_allow_html=True)
    st.markdown("<div class='small-muted' style='font-weight:600; margin-bottom:8px;'>Suggested next questions</div>", unsafe_allow_html=True)
    cols = st.columns(len(suggestions))
    for i, q in enumerate(suggestions):
        with cols[i]:
            if st.button(
                q,
                use_container_width=True,
                key=f"sugg_btn_{chat['id']}_{len(chat.get('messages', []))}_{i}",
            ):
                st.session_state["pending_prompt"] = q
                st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)


# =============================
# SIDEBAR
# =============================
with st.sidebar:
    st.markdown(f"### {APP_TITLE}")
    st.caption(APP_SUBTITLE)

    if st.button("➕ New analysis", use_container_width=True, key="btn_new_chat"):
        new_chat()
        st.rerun()

    st.divider()

    st.markdown("**Recent Analyses**")
    if not st.session_state.chats:
        st.caption("No analyses yet.")
    else:
        for c in st.session_state.chats[:12]:
            title = c.get("title") or "Untitled"
            if st.button(title, use_container_width=True, key=f"chat_select_{c['id']}"):
                _set_active(c["id"])
                st.rerun()

    st.divider()
    connected = is_api_alive()
    st.caption("Status: 🟢 Connected" if connected else "Status: 🔴 API Offline")


# =============================
# MAIN
# =============================
chat = get_chat()

st.markdown(
    f"""
<div class="right-header">
  <h2>{chat.get("title","New analysis")}</h2>
</div>
""",
    unsafe_allow_html=True,
)
st.caption("GARV returns business insight, calculation logic, detailed results, trend/comparison views, and admin trace.")


# =============================
# RENDER MESSAGES
# =============================
for idx, m in enumerate(chat["messages"]):
    role = m.get("role", "assistant")
    mtype = m.get("type", "text")
    payload = m.get("payload", {})
    content = (m.get("content") or "").strip()

    # USER bubble
    if role == "user":
        with st.chat_message("user"):
            if content:
                st.markdown(content)
        continue

    # ASSISTANT (NO bubble)
    with st.container():
        if mtype == "text":
            if content:
                st.markdown(content)
            continue

        if mtype != "answer":
            if content:
                st.markdown(content)
            continue

        out = payload if isinstance(payload, dict) else {}
        answer_card = out.get("answer_card") or {}
        final_sql = out.get("final_sql") or ""
        rows = out.get("rows") or []
        df = pd.DataFrame(rows) if isinstance(rows, list) and rows else pd.DataFrame()
        chart_spec = out.get("chart_spec") or out.get("chart") or None
        debug = out.get("debug") or {}
        trace = (debug.get("trace") or []) if isinstance(debug, dict) else []

        # Question (only if present)
        question_text = (out.get("question") or out.get("user_question") or "").strip()
        if question_text:
            st.markdown('<div class="card">', unsafe_allow_html=True)
            st.markdown('<div class="section-title">Your question</div>', unsafe_allow_html=True)
            st.markdown(f"“{question_text}”")
            st.markdown("</div>", unsafe_allow_html=True)
            st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # Executive Insight
        conf = (answer_card.get("confidence") or "MEDIUM").upper()
        badge_class = "badge-med"
        if conf == "HIGH":
            badge_class = "badge-high"
        elif conf == "LOW":
            badge_class = "badge-low"

        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown(
            f"""
<div style="display:flex;align-items:center;justify-content:space-between;gap:10px;">
  <div class="section-title">Executive Insight</div>
  <span class="badge {badge_class}">Confidence: {conf}</span>
</div>
""",
            unsafe_allow_html=True,
        )

        insight = (answer_card.get("answer") or "").strip()
        if not insight:
            exp = out.get("explanation") or {}
            if isinstance(exp, dict):
                insight = (exp.get("summary") or "").strip()

        st.markdown(insight or "Insight not available.")

        scope = answer_card.get("scope") or {}
        entities = scope.get("entities") or out.get("entities") or {}
        tables = scope.get("tables") or out.get("retrieved_tables") or []

        with st.expander("Details (scope & assumptions)", expanded=False):
            st.write({"intent": out.get("intent"), "entities": entities, "tables_used": tables})
            assumptions = out.get("assumptions") or []
            if assumptions:
                st.markdown("**Assumptions**")
                for a in assumptions:
                    st.markdown(f"- {a}")

        st.markdown("</div>", unsafe_allow_html=True)
        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # How we calculated it (SQL)
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">How we calculated it</div>', unsafe_allow_html=True)
        with st.expander("View calculation logic (SQL)", expanded=False):
            st.code(final_sql, language="sql")
            st.download_button(
                "Download SQL",
                data=final_sql.encode("utf-8"),
                file_name=f"{chat['id']}_sql_{idx}.sql",
                mime="text/plain",
                key=f"dl_sql_{chat['id']}_{idx}",
            )
        st.markdown("</div>", unsafe_allow_html=True)
        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # Detailed results
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Detailed results</div>', unsafe_allow_html=True)

        if df.empty:
            st.warning("No matching records for this filter.")
            sugg = out.get("suggested_questions") or []
            sugg = [str(s).strip() for s in sugg if str(s).strip()]

            if sugg:
                st.caption("Try one of these next:")
                cols = st.columns(min(3, len(sugg)))
                for i, s in enumerate(sugg[:3]):
                    with cols[i]:
                        if st.button(s, use_container_width=True, key=f"empty_sugg_{chat['id']}_{idx}_{i}"):
                            st.session_state["pending_prompt"] = s
                            st.rerun()
            else:
                st.caption("Tip: Try widening the time window or removing a filter.")
        else:
            st.caption(f"Rows returned: {len(df)}")
            st.dataframe(df, use_container_width=True, height=320)

            excel_bytes = df_to_excel_bytes(df)

            c1, c2 = st.columns([1, 1])
            with c1:
                st.download_button(
                    "⬇ Download CSV",
                    data=df.to_csv(index=False).encode("utf-8"),
                    file_name=f"{chat['id']}_results_{idx}.csv",
                    mime="text/csv",
                    key=f"dl_csv_{chat['id']}_{idx}",
                )
            with c2:
                st.download_button(
                    "⬇ Download Excel",
                    data=excel_bytes if excel_bytes else b"",
                    file_name=f"{chat['id']}_results_{idx}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    disabled=(excel_bytes is None),
                    key=f"dl_xlsx_{chat['id']}_{idx}",
                )
                if excel_bytes is None:
                    st.caption("Install openpyxl to enable Excel export: `pip install openpyxl`")

        st.markdown("</div>", unsafe_allow_html=True)
        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # Trend / comparison view
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Trend / comparison view</div>', unsafe_allow_html=True)

        overrides_for_chat = chat.get("viz_overrides", {})
        current_override = overrides_for_chat.get(str(idx), {}) if isinstance(overrides_for_chat, dict) else {}

        if df.empty:
            st.info("No data to visualize.")
        else:
            render_chart(df, chart_spec, override=current_override)

            with st.expander("Edit view", expanded=False):
                meta = _infer_columns(df)
                cols = meta["cols"]
                numeric = meta["numeric"]
                datetime_cols = meta["datetime"]
                categorical = meta["categorical"]

                default_type = (current_override.get("type") or (chart_spec or {}).get("type") or "bar")
                chart_type = st.selectbox(
                    "Chart type",
                    ["bar", "line", "scatter", "table"],
                    index=["bar", "line", "scatter", "table"].index(default_type)
                    if default_type in ["bar", "line", "scatter", "table"]
                    else 0,
                    key=f"viz_type_{chat['id']}_{idx}",
                )

                default_x = (
                    current_override.get("x")
                    or (chart_spec or {}).get("x")
                    or (datetime_cols[0] if datetime_cols else (categorical[0] if categorical else cols[0]))
                )
                default_y = current_override.get("y") or (chart_spec or {}).get("y") or (numeric[0] if numeric else None)

                x = st.selectbox(
                    "X axis",
                    cols,
                    index=cols.index(default_x) if default_x in cols else 0,
                    key=f"viz_x_{chat['id']}_{idx}",
                )

                y: Optional[str] = None
                if numeric:
                    y = st.selectbox(
                        "Y axis",
                        numeric,
                        index=numeric.index(default_y) if default_y in numeric else 0,
                        key=f"viz_y_{chat['id']}_{idx}",
                    )
                else:
                    st.caption("No numeric column detected for Y axis.")

                if st.button("Apply changes", use_container_width=True, key=f"apply_viz_{chat['id']}_{idx}"):
                    chat.setdefault("viz_overrides", {})
                    chat["viz_overrides"][str(idx)] = {"type": chart_type, "x": x, "y": y}
                    st.rerun()

        st.markdown("</div>", unsafe_allow_html=True)
        st.markdown('<div class="hr"></div>', unsafe_allow_html=True)

        # Actions (Export only)
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('<div class="section-title">Actions</div>', unsafe_allow_html=True)

        a1, a2 = st.columns([1, 1])
        with a1:
            pdf_bytes = build_pdf_bytes(out)
            st.download_button(
                "⬇ Export to PDF",
                data=pdf_bytes if pdf_bytes else b"",
                file_name=f"{chat['id']}_export_{idx}.pdf",
                mime="application/pdf",
                disabled=(pdf_bytes is None),
                use_container_width=True,
                key=f"pdf_{chat['id']}_{idx}",
            )
            if pdf_bytes is None:
                st.caption("Install reportlab to enable PDF export.")
        with a2:
            pptx_bytes = build_pptx_bytes(out)
            st.download_button(
                "⬇ Export to PPT",
                data=pptx_bytes if pptx_bytes else b"",
                file_name=f"{chat['id']}_export_{idx}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                disabled=(pptx_bytes is None),
                use_container_width=True,
                key=f"ppt_{chat['id']}_{idx}",
            )
            if pptx_bytes is None:
                st.caption("Install python-pptx to enable PPT export.")

        st.markdown("</div>", unsafe_allow_html=True)

        # Run trace (for admins)
        if trace:
            with st.expander("Run trace (for admins)", expanded=False):
                st.graphviz_chart(trace_to_dot(trace))
                st.json({"trace": trace})


# # =============================
# INPUT + RUN
# =============================

# ✅ Show suggestions right above the input (no navigation, no sticky overlay)
render_suggestions_inline(chat)

# ✅ Take prompt either from clicked suggestion or user input
pending = st.session_state.pop("pending_prompt", None)
prompt = pending or st.chat_input(
    "Ask a question… e.g. Top 5 airports by avg security wait time in last 30 days"
)

if prompt:
    if chat["title"] == "New analysis":
        chat["title"] = prompt[:34] + ("…" if len(prompt) > 34 else "")

    # store user message
    chat["messages"].append({"role": "user", "type": "text", "content": prompt})

    if not is_api_alive():
        chat["messages"].append(
            {"role": "assistant", "type": "text", "content": "❌ API is Offline. Start FastAPI first, then retry."}
        )
        st.rerun()

    status = st.status("Running analysis…", expanded=True)
    status.update(label="Understanding request…", state="running")
    time.sleep(0.10)
    status.update(label="Generating SQL & validating…", state="running")

    try:
        req_payload = {
            "question": prompt,
            "top_k_schema": 5,
            "return_rows": 50,
            "enable_viz": True,
            "chat_history": chat.get("chat_history", []),
            "memory_entities": chat.get("memory_entities", {}),
        }

        r = requests.post(API_BASE + TEXT2SQL_ENDPOINT, json=req_payload, timeout=180)
        out = r.json() if r is not None else {}

        if isinstance(out, dict):
            chat["chat_history"] = out.get("chat_history", chat.get("chat_history", []))
            chat["memory_entities"] = out.get("memory_entities", chat.get("memory_entities", {}))

        if isinstance(out, dict) and out.get("ok"):
            status.update(label="Completed", state="complete")
            out["question"] = prompt  # ✅ ensures "Your question" renders
            chat["messages"].append({"role": "assistant", "type": "answer", "payload": out})
            st.rerun()
        else:
            status.update(label="Failed", state="error")
            msg = out.get("message", "Unknown error") if isinstance(out, dict) else "Unknown error"
            chat["messages"].append({"role": "assistant", "type": "text", "content": f"❌ {msg}"})
            st.rerun()

    except Exception as e:
        status.update(label="Failed", state="error")
        chat["messages"].append({"role": "assistant", "type": "text", "content": f"❌ API error: {e}"})
        st.rerun()